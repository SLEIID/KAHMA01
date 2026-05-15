from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Kolory Kahma "Sapphire Navy" ──────────────────────────────────────────────
NAVY   = RGBColor(0x0c, 0x1e, 0x3c)
BLUE   = RGBColor(0x27, 0x61, 0xeb)
BLUE2  = RGBColor(0x1a, 0x4f, 0xc4)
LIGHT  = RGBColor(0xed, 0xf2, 0xfb)
WHITE  = RGBColor(0xff, 0xff, 0xff)
AMBER  = RGBColor(0xf5, 0x9e, 0x0b)
GREEN  = RGBColor(0x22, 0xc5, 0x5e)
RED    = RGBColor(0xef, 0x44, 0x44)
MUTED  = RGBColor(0x64, 0x74, 0x8b)
INK    = RGBColor(0x1e, 0x29, 0x3b)
CARD   = RGBColor(0xf8, 0xfa, 0xfc)

W  = Inches(13.33)   # widescreen 16:9
H  = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # całkowicie pusty

# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(5, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def add_pill(slide, text, x, y, bg, fg, size=11):
    """Rounded-rectangle badge."""
    tw = Cm(len(text) * 0.22 + 0.8)
    th = Cm(0.62)
    r = add_rect(slide, x, y, tw, th, fill=bg)
    r.adjustments[0] = 0.5
    tf = r.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = True
    run.font.color.rgb = fg
    return r, tw

def header_bar(slide, subtitle=""):
    """Granatowy pasek nagłówka."""
    add_rect(slide, 0, 0, W, Cm(1.85), fill=NAVY)
    # Logo kafelek
    logo = add_rect(slide, Cm(0.5), Cm(0.3), Cm(1.25), Cm(1.25), fill=BLUE)
    logo.adjustments[0] = 0.15
    add_text(slide, "K", Cm(0.5), Cm(0.28), Cm(1.25), Cm(1.3),
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "KAHMA", Cm(2.0), Cm(0.28), Cm(4), Cm(0.7),
             size=15, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, Cm(2.0), Cm(1.0), Cm(8), Cm(0.6),
                 size=10, color=RGBColor(0x8b, 0xa5, 0xd4))

def bullet_block(slide, items, x, y, w, icon_col=BLUE):
    """Lista kroków z numerowanymi kółkami."""
    cy = y
    step = Cm(1.1)
    for i, (title, desc) in enumerate(items):
        # kółko
        circ = add_rect(slide, x, cy, Cm(0.75), Cm(0.75), fill=icon_col)
        circ.adjustments[0] = 0.5
        tf = circ.text_frame
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(i+1)
        r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE
        # tekst
        add_text(slide, title, x + Cm(0.95), cy - Cm(0.05), w - Cm(1.0), Cm(0.5),
                 size=12, bold=True, color=INK)
        add_text(slide, desc, x + Cm(0.95), cy + Cm(0.42), w - Cm(1.0), Cm(0.5),
                 size=10, color=MUTED)
        cy += step
    return cy

def feat_row(slide, icon, title, desc, x, y, w):
    bg = add_rect(slide, x, y, w, Cm(0.9), fill=CARD, line_color=RGBColor(0xe2,0xe8,0xf0), line_w=Pt(0.75))
    bg.adjustments[0] = 0.1
    add_text(slide, icon, x + Cm(0.2), y + Cm(0.1), Cm(0.7), Cm(0.7), size=14, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + Cm(1.0), y + Cm(0.05), w - Cm(1.2), Cm(0.45),
             size=11, bold=True, color=INK)
    add_text(slide, desc, x + Cm(1.0), y + Cm(0.48), w - Cm(1.2), Cm(0.38),
             size=9.5, color=MUTED)

def section_label(slide, text, x, y, color=BLUE):
    add_text(slide, text, x, y, Cm(12), Cm(0.55),
             size=13, bold=True, color=color)

def info_box(slide, icon, text, x, y, w, bg, border, fg):
    box = add_rect(slide, x, y, w, Cm(0.9), fill=bg, line_color=border, line_w=Pt(1.5))
    box.adjustments[0] = 0.1
    add_text(slide, icon + "  " + text, x + Cm(0.3), y + Cm(0.12),
             w - Cm(0.5), Cm(0.7), size=10, color=fg)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 0 — OKŁADKA
# ══════════════════════════════════════════════════════════════════════════════
def slide_cover():
    sl = prs.slides.add_slide(BLANK)

    # Tło granatowe
    add_rect(sl, 0, 0, W, H, fill=NAVY)

    # Dekoracyjny pasek niebieski po lewej
    add_rect(sl, 0, 0, Cm(0.5), H, fill=BLUE)

    # Badge szkolenie
    badge = add_rect(sl, Cm(3.5), Cm(1.6), Cm(6.3), Cm(0.7),
                     fill=RGBColor(0x1a, 0x4f, 0xc4))
    badge.adjustments[0] = 0.5
    tf = badge.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "🎓  SZKOLENIE UŻYTKOWNIKA · 2026"
    r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = RGBColor(0x93, 0xbb, 0xff)

    # Tytuł
    add_text(sl, "Witaj w systemie", Cm(1.5), Cm(2.5), Cm(10.3), Cm(1.1),
             size=32, bold=False, color=RGBColor(0xa8,0xc4,0xf0), align=PP_ALIGN.CENTER)
    add_text(sl, "KAHMA", Cm(1.5), Cm(3.4), Cm(10.3), Cm(2.0),
             size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(sl, "Twój kompletny system Mini ERP — prosto i mobilnie",
             Cm(1.5), Cm(5.35), Cm(10.3), Cm(0.7),
             size=16, color=RGBColor(0xa8,0xc4,0xf0), align=PP_ALIGN.CENTER)

    # Moduły — chipsy
    modules = ["1  Platforma ERP","2  Logowanie","3  HR i urlopy",
               "4  Wypożyczalnia","5  Materiałówka","6  Raport dnia","7  Rozliczenie"]
    cols = 4
    bw, bh = Cm(3.4), Cm(0.62)
    gap_x, gap_y = Cm(0.25), Cm(0.2)
    start_x = (W - (cols * bw + (cols-1) * gap_x)) / 2
    for i, m in enumerate(modules):
        col = i % cols; row = i // cols
        bx = start_x + col * (bw + gap_x)
        by = Cm(6.1) + row * (bh + gap_y)
        b = add_rect(sl, bx, by, bw, bh,
                     fill=RGBColor(0x10,0x25,0x50))
        b.adjustments[0] = 0.5
        tf = b.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = m
        r.font.size = Pt(11); r.font.bold = True
        r.font.color.rgb = RGBColor(0xd0,0xe4,0xff)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — PLATFORMA MINI ERP
# ══════════════════════════════════════════════════════════════════════════════
def slide_erp():
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, fill=LIGHT)
    header_bar(sl, "Mini ERP")

    # Numer slajdu
    add_text(sl, "1 / 7", W - Cm(2.5), Cm(0.55), Cm(2), Cm(0.6),
             size=10, color=RGBColor(0x8b,0xa5,0xd4), align=PP_ALIGN.RIGHT)

    # Tytuł
    add_text(sl, "🏢  Platforma Mini ERP", Cm(0.7), Cm(2.1), Cm(12), Cm(0.8),
             size=26, bold=True, color=NAVY)
    add_text(sl, "Jeden system — wszystko czego potrzebujesz w pracy",
             Cm(0.7), Cm(2.95), Cm(12), Cm(0.5), size=12, color=MUTED)

    # Moduły — siatka 3×2
    mods = [
        ("📋","Raporty Dnia","Czas, lokalizacje, materiały"),
        ("🔧","Wypożyczalnia","Sprzęt, narzędzia, podnośniki"),
        ("📦","Materiałówka","4 535 materiałów z wyszukiwarką"),
        ("🏖️","HR i Urlopy","Wnioski, 26 dni, kalendarze"),
        ("🚗","Flota Pojazdów","21 pojazdów, liczniki km"),
        ("📊","Eksport XLSX","Raporty jednym kliknięciem"),
    ]
    cols = 3; cw = Cm(3.7); ch = Cm(1.6)
    gx = Cm(0.3); gy = Cm(0.3)
    ox = Cm(0.7); oy = Cm(3.65)
    for i, (icon, title, desc) in enumerate(mods):
        col = i % cols; row = i // cols
        bx = ox + col * (cw + gx); by = oy + row * (ch + gy)
        box = add_rect(sl, bx, by, cw, ch, fill=WHITE,
                       line_color=RGBColor(0xe2,0xe8,0xf0), line_w=Pt(1))
        box.adjustments[0] = 0.1
        add_text(sl, icon, bx + Cm(0.3), by + Cm(0.22), Cm(0.8), Cm(0.8), size=22, align=PP_ALIGN.CENTER)
        add_text(sl, title, bx + Cm(1.15), by + Cm(0.2), cw - Cm(1.3), Cm(0.55),
                 size=12, bold=True, color=NAVY)
        add_text(sl, desc, bx + Cm(1.15), by + Cm(0.78), cw - Cm(1.3), Cm(0.65),
                 size=9.5, color=MUTED)

    # Info box
    info_box(sl, "💡",
             "Działa w przeglądarce na telefonie i komputerze — kahma.leanmatik.net",
             Cm(0.7), Cm(6.78), Cm(11.9),
             RGBColor(0xef,0xf6,0xff), RGBColor(0xbf,0xdb,0xfe), RGBColor(0x1d,0x40,0xaf))

    # Role — prawa kolumna
    rx = Cm(12.6)
    add_text(sl, "Dwie role", rx, Cm(2.1), Cm(4.5), Cm(0.6),
             size=14, bold=True, color=NAVY)

    # Admin
    adm = add_rect(sl, rx, Cm(2.75), Cm(4.5), Cm(2.2),
                   fill=RGBColor(0xef,0xf6,0xff), line_color=RGBColor(0xbf,0xdb,0xfe), line_w=Pt(1.5))
    adm.adjustments[0] = 0.1
    add_text(sl, "🛡️  Administrator", rx+Cm(0.2), Cm(2.8), Cm(4.1), Cm(0.5),
             size=12, bold=True, color=NAVY)
    adm_items = ["Raporty wszystkich","Zarządza flotą i użytk.","Zatwierdza urlopy","Eksport XLSX"]
    for j, t in enumerate(adm_items):
        add_text(sl, "•  "+t, rx+Cm(0.3), Cm(3.32)+j*Cm(0.38), Cm(4.1), Cm(0.38),
                 size=10, color=INK)

    # Pracownik
    prc = add_rect(sl, rx, Cm(5.1), Cm(4.5), Cm(1.95),
                   fill=RGBColor(0xf0,0xfd,0xf4), line_color=RGBColor(0xbb,0xf7,0xd0), line_w=Pt(1.5))
    prc.adjustments[0] = 0.1
    add_text(sl, "👷  Pracownik", rx+Cm(0.2), Cm(5.15), Cm(4.1), Cm(0.5),
             size=12, bold=True, color=NAVY)
    prc_items = ["Tworzy własne raporty","Wypożycza sprzęt","Rejestruje materiały","Wnioski urlopowe"]
    for j, t in enumerate(prc_items):
        add_text(sl, "•  "+t, rx+Cm(0.3), Cm(5.67)+j*Cm(0.37), Cm(4.1), Cm(0.37),
                 size=10, color=INK)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — LOGOWANIE
# ══════════════════════════════════════════════════════════════════════════════
def slide_login():
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, fill=LIGHT)
    header_bar(sl, "Logowanie do systemu")
    add_text(sl, "2 / 7", W-Cm(2.5), Cm(0.55), Cm(2), Cm(0.6),
             size=10, color=RGBColor(0x8b,0xa5,0xd4), align=PP_ALIGN.RIGHT)

    add_text(sl, "🔐  Logowanie do systemu", Cm(0.7), Cm(2.1), Cm(12), Cm(0.8),
             size=26, bold=True, color=NAVY)
    add_text(sl, "Bezpieczny dostęp przez przeglądarkę — telefon lub komputer",
             Cm(0.7), Cm(2.95), Cm(12), Cm(0.5), size=12, color=MUTED)

    # Kroki
    bullet_block(sl, [
        ("Otwórz przeglądarkę", "Wejdź na kahma.leanmatik.net — działa na każdym urządzeniu"),
        ("Wpisz login i hasło",  "Dane dostępowe otrzymujesz od administratora systemu"),
        ("Kliknij <<Zaloguj się>>","Przekierowanie na Dashboard z dostępnymi modułami"),
        ("Wylogowanie",          "Przycisk w lewym dolnym rogu paska bocznego"),
    ], Cm(0.7), Cm(3.6), Cm(8.5))

    info_box(sl, "⚠️",
             "Po 15 min bezczynności sesja odświeża się automatycznie. Dane w formularzu nie przepadają.",
             Cm(0.7), Cm(6.8), Cm(8.5),
             RGBColor(0xff,0xfb,0xeb), RGBColor(0xfd,0xe6,0x8a), RGBColor(0x92,0x40,0x0e))

    # Mockup ekranu logowania
    mx = Cm(9.7); mw = Cm(7.3); my = Cm(2.3)
    bg = add_rect(sl, mx, my, mw, Cm(4.7), fill=RGBColor(0x27,0x27,0x2a))
    bg.adjustments[0] = 0.1
    # logo kahma w mockupie
    klogo = add_rect(sl, mx + Cm(2.0), my + Cm(0.4), Cm(3.3), Cm(0.9),
                     fill=BLUE)
    klogo.adjustments[0] = 0.12
    tf = klogo.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "KAHMA"
    r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = WHITE

    add_text(sl, "Login", mx+Cm(0.4), my+Cm(1.5), Cm(2), Cm(0.3),
             size=9, bold=True, color=RGBColor(0x8b,0xa5,0xd4))
    inp1 = add_rect(sl, mx+Cm(0.4), my+Cm(1.82), mw-Cm(0.8), Cm(0.55),
                    fill=RGBColor(0x3f,0x3f,0x46), line_color=RGBColor(0x52,0x52,0x5b), line_w=Pt(1))
    inp1.adjustments[0] = 0.1
    add_text(sl, "jan.kowalski", mx+Cm(0.7), my+Cm(1.88), mw-Cm(1.4), Cm(0.4),
             size=11, color=WHITE)

    add_text(sl, "Hasło", mx+Cm(0.4), my+Cm(2.5), Cm(2), Cm(0.3),
             size=9, bold=True, color=RGBColor(0x8b,0xa5,0xd4))
    inp2 = add_rect(sl, mx+Cm(0.4), my+Cm(2.82), mw-Cm(0.8), Cm(0.55),
                    fill=RGBColor(0x3f,0x3f,0x46), line_color=RGBColor(0x52,0x52,0x5b), line_w=Pt(1))
    inp2.adjustments[0] = 0.1
    add_text(sl, "••••••••", mx+Cm(0.7), my+Cm(2.88), mw-Cm(1.4), Cm(0.4),
             size=11, color=RGBColor(0x94,0xa3,0xb8))

    btn = add_rect(sl, mx+Cm(0.4), my+Cm(3.55), mw-Cm(0.8), Cm(0.7), fill=BLUE)
    btn.adjustments[0] = 0.1
    tf = btn.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "🔐  Zaloguj się"
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE

    # Po zalogowaniu
    add_text(sl, "Po zalogowaniu — Dashboard", Cm(9.7), Cm(7.2), Cm(7.3), Cm(0.5),
             size=12, bold=True, color=NAVY)
    dash_items = [("📊","Twoje statystyki miesiąca"),("🔔","Przypomnienie o raporcie po 12:00"),("⬛","Kafelki modułów")]
    for j, (ic, txt) in enumerate(dash_items):
        add_text(sl, ic+"  "+txt, Cm(9.9), Cm(7.7)+j*Cm(0.45), Cm(7), Cm(0.42),
                 size=10.5, color=INK)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — HR
# ══════════════════════════════════════════════════════════════════════════════
def slide_hr():
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, fill=LIGHT)
    header_bar(sl, "Moduł HR i Urlopy")
    add_text(sl, "3 / 7", W-Cm(2.5), Cm(0.55), Cm(2), Cm(0.6),
             size=10, color=RGBColor(0x8b,0xa5,0xd4), align=PP_ALIGN.RIGHT)

    add_text(sl, "🏖️  Moduł HR — Urlopy i Obecność", Cm(0.7), Cm(2.1), Cm(12), Cm(0.8),
             size=26, bold=True, color=NAVY)
    add_text(sl, "Wnioski urlopowe, saldo dni i kalendarz pracownika",
             Cm(0.7), Cm(2.95), Cm(12), Cm(0.5), size=12, color=MUTED)

    # Lewa kolumna — typy urlopów
    add_text(sl, "Typy urlopów", Cm(0.7), Cm(3.6), Cm(6), Cm(0.5),
             size=13, bold=True, color=NAVY)

    urlopers = [
        (RGBColor(0xff,0xfb,0xeb), RGBColor(0xfd,0xe6,0x8a), "🌴  Urlop wypoczynkowy", "odlicza z puli 26 dni"),
        (RGBColor(0xff,0xfb,0xeb), RGBColor(0xfd,0xe6,0x8a), "⚡  Urlop na żądanie",    "odlicza z puli"),
        (RGBColor(0xff,0xfb,0xeb), RGBColor(0xfd,0xe6,0x8a), "🎉  Urlop okolicznościowy","odlicza z puli"),
        (RGBColor(0xf0,0xfd,0xf4), RGBColor(0xbb,0xf7,0xd0), "🏥  L4 — Zwolnienie",    "NIE odlicza z puli"),
        (RGBColor(0xf1,0xf5,0xf9), RGBColor(0xe2,0xe8,0xf0), "💼  Urlop bezpłatny",     "odlicza z puli"),
    ]
    for i, (bg, bd, name, note) in enumerate(urlopers):
        bx = Cm(0.7); by = Cm(4.12) + i * Cm(0.62)
        box = add_rect(sl, bx, by, Cm(8.2), Cm(0.55), fill=bg,
                       line_color=bd, line_w=Pt(1))
        box.adjustments[0] = 0.1
        add_text(sl, name, bx+Cm(0.3), by+Cm(0.06), Cm(5.5), Cm(0.42),
                 size=11, bold=True, color=INK)
        fg = RGBColor(0x15,0x80,0x3d) if "NIE" in note else RGBColor(0xb4,0x53,0x09)
        nbg = RGBColor(0xdc,0xfc,0xe7) if "NIE" in note else RGBColor(0xfe,0xf3,0xc7)
        pill = add_rect(sl, bx+Cm(5.9), by+Cm(0.1), Cm(2.15), Cm(0.36),
                        fill=nbg)
        pill.adjustments[0] = 0.5
        tf = pill.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = note
        r.font.size = Pt(9); r.font.bold = True; r.font.color.rgb = fg

    # Kroki złożenia wniosku
    add_text(sl, "Jak złożyć wniosek?", Cm(0.7), Cm(7.35), Cm(8.2), Cm(0.5),
             size=12, bold=True, color=NAVY)
    steps_h = [
        ("HR → <<Nowy wniosek>>", "Wybierz typ urlopu i daty"),
        ("Dodaj uwagi",          "Opcjonalne uzasadnienie"),
        ("Czekaj na decyzję",    "Admin zatwierdza lub odrzuca"),
    ]
    for i, (t, d) in enumerate(steps_h):
        circ = add_rect(sl, Cm(0.7)+i*Cm(2.75), Cm(7.85), Cm(0.6), Cm(0.6), fill=BLUE)
        circ.adjustments[0] = 0.5
        tf = circ.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(i+1)
        r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE
        add_text(sl, t, Cm(1.45)+i*Cm(2.75), Cm(7.85), Cm(2.5), Cm(0.35),
                 size=10, bold=True, color=INK)
        add_text(sl, d, Cm(1.45)+i*Cm(2.75), Cm(8.22), Cm(2.5), Cm(0.35),
                 size=9, color=MUTED)

    # Prawa — saldo + kalendarz
    rx = Cm(9.4)
    add_text(sl, "Saldo urlopowe", rx, Cm(3.6), Cm(5.5), Cm(0.5),
             size=13, bold=True, color=NAVY)

    saldo = [("26","Dni w roku",RGBColor(0xf0,0xfd,0xf4),RGBColor(0xbb,0xf7,0xd0),RGBColor(0x16,0xa3,0x4a)),
             ("14","Wykorzystane",RGBColor(0xef,0xf6,0xff),RGBColor(0xbf,0xdb,0xfe),BLUE),
             ("12","Pozostałe",RGBColor(0xff,0xfb,0xeb),RGBColor(0xfd,0xe6,0x8a),RGBColor(0xd9,0x77,0x06))]
    for i, (num, lbl, bg, bd, fg) in enumerate(saldo):
        bx = rx + i * Cm(1.75)
        box = add_rect(sl, bx, Cm(4.15), Cm(1.6), Cm(1.3), fill=bg,
                       line_color=bd, line_w=Pt(1.5))
        box.adjustments[0] = 0.1
        add_text(sl, num, bx, Cm(4.2), Cm(1.6), Cm(0.75),
                 size=28, bold=True, color=fg, align=PP_ALIGN.CENTER)
        add_text(sl, lbl, bx, Cm(4.95), Cm(1.6), Cm(0.42),
                 size=9, bold=True, color=fg, align=PP_ALIGN.CENTER)

    # Mini kalendarz
    add_text(sl, "Kalendarz obecności — kwiecień", rx, Cm(5.65), Cm(5.5), Cm(0.5),
             size=11, bold=True, color=NAVY)

    days_header = ["Pn","Wt","Śr","Cz","Pt","So","Nd"]
    cal_data = [
        [None,True,True,True,True,False,False],
        [True,True,True,True,True,False,False],
        ["L","L","L",True,True,False,False],
        ["T",None,None,None,None,None,None],
    ]
    cs = Cm(0.62); cgx = Cm(0.08)
    cal_ox = rx; cal_oy = Cm(6.15)
    for c, d in enumerate(days_header):
        add_text(sl, d, cal_ox + c*(cs+cgx), cal_oy, cs, Cm(0.32),
                 size=8, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    day_num = 1
    for row in range(4):
        for col in range(7):
            val = cal_data[row][col] if row < len(cal_data) and col < len(cal_data[row]) else None
            bx2 = cal_ox + col*(cs+cgx)
            by2 = cal_oy + Cm(0.38) + row*(cs+cgx)
            if val is None and row == 0 and col == 0:
                day_num += 1; continue
            if val is None: continue
            if val == "T":
                bg2 = BLUE; fg2 = WHITE
            elif val == "L":
                bg2 = RGBColor(0xdb,0xea,0xfe); fg2 = RGBColor(0x1d,0x4e,0xd8)
            elif val is True:
                bg2 = RGBColor(0xdc,0xfc,0xe7); fg2 = RGBColor(0x15,0x80,0x3d)
            else:
                bg2 = RGBColor(0xf1,0xf5,0xf9); fg2 = MUTED
            box2 = add_rect(sl, bx2, by2, cs, cs, fill=bg2)
            box2.adjustments[0] = 0.1
            add_text(sl, str(day_num), bx2, by2+Cm(0.1), cs, cs-Cm(0.1),
                     size=9, bold=True, color=fg2, align=PP_ALIGN.CENTER)
            day_num += 1

    # legenda
    legend = [
        (RGBColor(0xdc,0xfc,0xe7), "Przepracowany"),
        (RGBColor(0xdb,0xea,0xfe), "Urlop"),
    ]
    for i, (dot_color, lbl) in enumerate(legend):
        bx3 = rx + i*Cm(2.5)
        dot = add_rect(sl, bx3, Cm(7.92), Cm(0.32), Cm(0.32), fill=dot_color)
        dot.adjustments[0] = 0.1
        add_text(sl, lbl, bx3+Cm(0.38), Cm(7.9), Cm(2.0), Cm(0.38), size=9, color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — WYPOŻYCZALNIA
# ══════════════════════════════════════════════════════════════════════════════
def slide_rental():
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, fill=LIGHT)
    header_bar(sl, "Wypożyczalnia Sprzętu")
    add_text(sl, "4 / 7", W-Cm(2.5), Cm(0.55), Cm(2), Cm(0.6),
             size=10, color=RGBColor(0x8b,0xa5,0xd4), align=PP_ALIGN.RIGHT)

    add_text(sl, "🔧  Wypożyczalnia Sprzętu", Cm(0.7), Cm(2.1), Cm(12), Cm(0.8),
             size=26, bold=True, color=NAVY)
    add_text(sl, "Narzędzia, podnośniki i sprzęt montażowy — szybkie wypożyczenie i zwrot",
             Cm(0.7), Cm(2.95), Cm(12), Cm(0.5), size=12, color=MUTED)

    bullet_block(sl, [
        ("Przejdź do <<Wypożyczalnia>>",    "Lista dostępnego sprzętu posortowana wg kategorii"),
        ("Znajdź potrzebny sprzęt",        "Filtruj po kategorii — status ZIELONY = dostępny"),
        ("Kliknij <<Wypożycz>>",            "Sprzęt przypisany do Ciebie, status: WYPOŻYCZONY"),
        ("Zwróć sprzęt po pracy",          "Twoje wypożyczenia → <<Zwróć>> → opcjonalna notatka"),
    ], Cm(0.7), Cm(3.6), Cm(8.5))

    add_text(sl, "Zgłoś problem ze sprzętem", Cm(0.7), Cm(6.2), Cm(8.5), Cm(0.5),
             size=12, bold=True, color=NAVY)
    feat_row(sl, "⚠️", "Tylko dla wypożyczonego przez Ciebie sprzętu",
             "Przycisk <<Problem>> widoczny przy Twoich aktywnych wypożyczeniach",
             Cm(0.7), Cm(6.75), Cm(8.5))
    feat_row(sl, "📝", "Opisz usterkę",
             "Admin otrzyma zgłoszenie i zamknie je po naprawie",
             Cm(0.7), Cm(7.73), Cm(8.5))

    # Statusy
    rx = Cm(9.5)
    add_text(sl, "Statusy sprzętu", rx, Cm(3.6), Cm(4.5), Cm(0.5),
             size=13, bold=True, color=NAVY)
    statuses = [
        (RGBColor(0xf0,0xfd,0xf4), RGBColor(0xbb,0xf7,0xd0), RGBColor(0x22,0xc5,0x5e), "Dostępny",    "Można wypożyczyć"),
        (RGBColor(0xef,0xf6,0xff), RGBColor(0xbf,0xdb,0xfe), BLUE,                       "Wypożyczony", "Zajęty przez innego"),
        (RGBColor(0xff,0xfb,0xeb), RGBColor(0xfd,0xe6,0x8a), AMBER,                      "Serwis",      "W naprawie — niedostępny"),
        (RGBColor(0xf1,0xf5,0xf9), RGBColor(0xe2,0xe8,0xf0), MUTED,                      "Wycofany",    "Nie pojawia się na liście"),
    ]
    for i, (bg, bd, dot_c, title, desc) in enumerate(statuses):
        by2 = Cm(4.2) + i * Cm(0.85)
        box2 = add_rect(sl, rx, by2, Cm(4.5), Cm(0.78), fill=bg,
                        line_color=bd, line_w=Pt(1))
        box2.adjustments[0] = 0.1
        dot2 = add_rect(sl, rx+Cm(0.25), by2+Cm(0.28), Cm(0.26), Cm(0.26), fill=dot_c)
        dot2.adjustments[0] = 0.5
        add_text(sl, title, rx+Cm(0.65), by2+Cm(0.08), Cm(3.6), Cm(0.38),
                 size=12, bold=True, color=INK)
        add_text(sl, desc, rx+Cm(0.65), by2+Cm(0.44), Cm(3.6), Cm(0.3),
                 size=9.5, color=MUTED)

    add_text(sl, "Dostępne kategorie", rx, Cm(7.6), Cm(4.5), Cm(0.5),
             size=12, bold=True, color=NAVY)
    cats = ["🏗️  Podnośniki","🔩  Sprzęt montażowy","🔦  Pomiarowy","⚡  Elektronarzędzia"]
    for i, c in enumerate(cats):
        pill2 = add_rect(sl, rx + (i%2)*Cm(2.2), Cm(8.12)+(i//2)*Cm(0.52), Cm(2.1), Cm(0.44),
                         fill=CARD, line_color=RGBColor(0xe2,0xe8,0xf0), line_w=Pt(1))
        pill2.adjustments[0] = 0.5
        add_text(sl, c, rx+(i%2)*Cm(2.2), Cm(8.14)+(i//2)*Cm(0.52), Cm(2.1), Cm(0.4),
                 size=9.5, color=INK, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — MATERIAŁÓWKA
# ══════════════════════════════════════════════════════════════════════════════
def slide_materials():
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, fill=LIGHT)
    header_bar(sl, "Materiałówka")
    add_text(sl, "5 / 7", W-Cm(2.5), Cm(0.55), Cm(2), Cm(0.6),
             size=10, color=RGBColor(0x8b,0xa5,0xd4), align=PP_ALIGN.RIGHT)

    add_text(sl, "📦  Materiałówka", Cm(0.7), Cm(2.1), Cm(12), Cm(0.8),
             size=26, bold=True, color=NAVY)
    add_text(sl, "4 535 materiałów — wyszukiwanie, rejestracja pobrań i alerty niskiego stanu",
             Cm(0.7), Cm(2.95), Cm(12), Cm(0.5), size=12, color=MUTED)

    bullet_block(sl, [
        ("Wejdź do <<Materiałówka>>",       "Widoczne od razu 30 ostatnio używanych materiałów"),
        ("Wyszukaj materiał (min. 3 znaki)","Kolejność słów nie ma znaczenia — system znajdzie po każdym"),
        ("Wybierz materiał z listy",       "Kliknij na pozycję aby ją wybrać"),
        ("Podaj ilość i jednostkę",        "Jednostki: szt / mb / kg / kpl / rolka / opak / l"),
        ("Kliknij <<Pobierz>>",             "Pobranie w historii — można usunąć tylko w dniu pobrania"),
    ], Cm(0.7), Cm(3.6), Cm(8.5))

    # Alert
    rx = Cm(9.5)
    add_text(sl, "Alert niskiego stanu", rx, Cm(3.6), Cm(4.5), Cm(0.5),
             size=13, bold=True, color=NAVY)
    alert_items = [
        ("📷","Zaznacz <<Zgłoś niski stan>>","Checkbox przy rejestracji pobrania"),
        ("🖼️","Dodaj zdjęcie (opcja)",       "Z kamery lub galerii telefonu"),
        ("🔔","Admin widzi alert",            "Zamyka po uzupełnieniu magazynu"),
    ]
    for i, (ic, t, d) in enumerate(alert_items):
        feat_row(sl, ic, t, d, rx, Cm(4.2)+i*Cm(1.0), Cm(4.5))

    add_text(sl, "Materiały w raporcie dnia", rx, Cm(7.35), Cm(4.5), Cm(0.5),
             size=12, bold=True, color=NAVY)

    info_box(sl, "💡",
             "Dodawaj materiały bezpośrednio do wpisu raportu — lokalizacja wypełnia się automatycznie.",
             rx, Cm(7.88), Cm(4.5),
             RGBColor(0xef,0xf6,0xff), RGBColor(0xbf,0xdb,0xfe), RGBColor(0x1d,0x40,0xaf))

    # Przykładowe pobranie
    ex = add_rect(sl, Cm(0.7), Cm(7.1), Cm(8.5), Cm(1.1),
                  fill=WHITE, line_color=RGBColor(0xe2,0xe8,0xf0), line_w=Pt(1))
    ex.adjustments[0] = 0.1
    add_text(sl, "Przykładowe pobranie", Cm(1.0), Cm(7.15), Cm(8), Cm(0.35),
             size=9, color=MUTED)
    add_text(sl, "Kabel YDY 3×1,5mm²  —  15 mb", Cm(1.0), Cm(7.5), Cm(8), Cm(0.45),
             size=13, bold=True, color=INK)
    add_text(sl, "Budowa A  ·  Hala produkcyjna  ·  dziś 09:42",
             Cm(1.0), Cm(7.95), Cm(8), Cm(0.35), size=10, color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — RAPORT DNIA
# ══════════════════════════════════════════════════════════════════════════════
def slide_report():
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, fill=LIGHT)
    header_bar(sl, "Raport Dnia")
    add_text(sl, "6 / 7", W-Cm(2.5), Cm(0.55), Cm(2), Cm(0.6),
             size=10, color=RGBColor(0x8b,0xa5,0xd4), align=PP_ALIGN.RIGHT)

    add_text(sl, "📋  Nowy Raport · Edycja · Podpis", Cm(0.7), Cm(2.1), Cm(12), Cm(0.8),
             size=24, bold=True, color=NAVY)
    add_text(sl, "Jeden raport na dzień — wiele wpisów, wiele lokalizacji, podpisy ekipy",
             Cm(0.7), Cm(2.95), Cm(12), Cm(0.5), size=12, color=MUTED)

    bullet_block(sl, [
        ("<<Raporty>> → <<Nowy raport>>",     "Data ustawiana automatycznie na dziś — niezmieniana"),
        ("Wybierz ekipę (raz na dzień)",   "Ekipa A/B/C/D/E/F — po wyborze zablokowana"),
        ("Dodaj wpis — <<+ Nowy wpis>>",    "Godziny od-do, lokalizacja, wydział, opis, pojazdy"),
        ("Krok 2 wpisu — materiały",       "Po zapisaniu wpisu dodaj materiały lub pomiń"),
        ("Możesz dodać wiele wpisów",      "Np. rano hala, popołudniu budowa — każdy wpis osobno"),
    ], Cm(0.7), Cm(3.6), Cm(8.5))

    info_box(sl, "🔒",
             "Edycja tylko do północy. Raporty z poprzednich dni zablokowane. Poproś admina o odblokowanie na 24h.",
             Cm(0.7), Cm(7.25), Cm(8.5),
             RGBColor(0xff,0xfb,0xeb), RGBColor(0xfd,0xe6,0x8a), RGBColor(0x92,0x40,0x0e))

    # Pojazdy
    rx = Cm(9.5)
    add_text(sl, "Pojazdy we wpisie", rx, Cm(3.6), Cm(4.5), Cm(0.5),
             size=13, bold=True, color=NAVY)

    vbox = add_rect(sl, rx, Cm(4.15), Cm(4.5), Cm(1.5),
                    fill=WHITE, line_color=RGBColor(0xe2,0xe8,0xf0), line_w=Pt(1))
    vbox.adjustments[0] = 0.1
    add_text(sl, "🚗  WGS 12345  ·  120 km", rx+Cm(0.3), Cm(4.28), Cm(4.1), Cm(0.5),
             size=12, bold=True, color=INK)
    add_text(sl, "🚐  WGS 98765  ·  85 km",  rx+Cm(0.3), Cm(4.82), Cm(4.1), Cm(0.5),
             size=12, bold=True, color=INK)
    add_text(sl, "Jeden wpis = wiele pojazdów. Liczniki km aktualizowane automatycznie.",
             rx, Cm(5.72), Cm(4.5), Cm(0.55), size=9.5, color=MUTED)

    # Podpis
    add_text(sl, "Podpisywanie raportu kolegi", rx, Cm(6.4), Cm(4.5), Cm(0.5),
             size=13, bold=True, color=NAVY)
    sign_items = [
        ("✍️","Przycisk <<Podpisz się>>",       "Widoczny na górze ekranu"),
        ("👥","Raporty kolegów z dzisiaj",      "Tylko bieżący dzień"),
        ("📝","Pełny dostęp do edycji",         "Dodawaj i edytuj wpisy"),
        ("↩️","Cofnięcie podpisu — tylko dziś","Przycisk <<Cofnij>> na karcie"),
    ]
    for i, (ic, t, d) in enumerate(sign_items):
        feat_row(sl, ic, t, d, rx, Cm(6.95)+i*Cm(1.0), Cm(4.5))

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — ROZLICZENIE DZIENNE
# ══════════════════════════════════════════════════════════════════════════════
def slide_settlement():
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, W, H, fill=LIGHT)
    header_bar(sl, "Rozliczenie dzienne")
    add_text(sl, "7 / 7", W-Cm(2.5), Cm(0.55), Cm(2), Cm(0.6),
             size=10, color=RGBColor(0x8b,0xa5,0xd4), align=PP_ALIGN.RIGHT)

    add_text(sl, "📊  Rozliczenie Dzienne", Cm(0.7), Cm(2.1), Cm(12), Cm(0.8),
             size=26, bold=True, color=NAVY)
    add_text(sl, "Podgląd pracy całej ekipy, zatwierdzanie i eksport do Excela",
             Cm(0.7), Cm(2.95), Cm(12), Cm(0.5), size=12, color=MUTED)

    # Tabela
    add_text(sl, "Widok admina — lista raportów z filtrami", Cm(0.7), Cm(3.55), Cm(9), Cm(0.5),
             size=12, bold=True, color=NAVY)

    # Filtry mockup
    filters_box = add_rect(sl, Cm(0.7), Cm(4.1), Cm(8.6), Cm(1.15),
                           fill=WHITE, line_color=RGBColor(0xe2,0xe8,0xf0), line_w=Pt(1))
    filters_box.adjustments[0] = 0.1
    filters = [("Od daty","2026-04-14"),("Do daty","2026-04-20"),("Pracownik","Wszyscy"),("Lokalizacja","Wszystkie")]
    for i, (lbl, val) in enumerate(filters):
        fx = Cm(0.95) + i * Cm(2.1)
        add_text(sl, lbl, fx, Cm(4.15), Cm(2), Cm(0.28), size=8, bold=True, color=MUTED)
        finp = add_rect(sl, fx, Cm(4.44), Cm(1.95), Cm(0.55),
                        fill=RGBColor(0xf1,0xf5,0xf9),
                        line_color=RGBColor(0xe2,0xe8,0xf0), line_w=Pt(1))
        finp.adjustments[0] = 0.08
        add_text(sl, val, fx+Cm(0.15), Cm(4.5), Cm(1.8), Cm(0.42), size=10, color=INK)

    # Nagłówek tabeli
    cols_t = [("Pracownik",Cm(2.2)),("Data",Cm(1.1)),("Ekipa",Cm(1.4)),
              ("Godziny",Cm(2.0)),("Wpisy",Cm(1.0)),("Status",Cm(1.8))]
    tx = Cm(0.7); ty = Cm(5.35)
    for cname, cw2 in cols_t:
        hdr = add_rect(sl, tx, ty, cw2, Cm(0.5), fill=NAVY)
        if tx == Cm(0.7): hdr.adjustments[0] = 0.0
        add_text(sl, cname, tx+Cm(0.1), ty+Cm(0.07), cw2-Cm(0.1), Cm(0.38),
                 size=9.5, bold=True, color=WHITE)
        tx += cw2

    rows = [
        ("Jan Kowalski","20.04","Ekipa A","07:00–15:30","2","Oczekuje"),
        ("Marek Nowak","20.04","Ekipa B","06:00–14:00","3","Zatwierdzone"),
        ("Piotr Wiśniewski","19.04","Ekipa C","08:00–16:00","1","Zatwierdzone"),
    ]
    row_bgs = [WHITE, RGBColor(0xf8,0xfa,0xfc), WHITE]
    for ri, row in enumerate(rows):
        tx2 = Cm(0.7); ry = Cm(5.85) + ri*Cm(0.62)
        for ci, (val, cw2) in enumerate(zip(row, [cw for _,cw in cols_t])):
            cell = add_rect(sl, tx2, ry, cw2, Cm(0.58), fill=row_bgs[ri],
                            line_color=RGBColor(0xe2,0xe8,0xf0), line_w=Pt(0.5))
            if ci == 5:
                pill_bg = RGBColor(0xdc,0xfc,0xe7) if val=="Zatwierdzone" else RGBColor(0xfe,0xf3,0xc7)
                pill_fg = RGBColor(0x15,0x80,0x3d) if val=="Zatwierdzone" else RGBColor(0xb4,0x53,0x09)
                pb = add_rect(sl, tx2+Cm(0.1), ry+Cm(0.11), cw2-Cm(0.2), Cm(0.36),
                              fill=pill_bg)
                pb.adjustments[0] = 0.5
                add_text(sl, val, tx2+Cm(0.1), ry+Cm(0.13), cw2-Cm(0.2), Cm(0.34),
                         size=9, bold=True, color=pill_fg, align=PP_ALIGN.CENTER)
            else:
                add_text(sl, val, tx2+Cm(0.1), ry+Cm(0.1), cw2-Cm(0.15), Cm(0.42),
                         size=10, color=INK)
            tx2 += cw2

    # Prawa kolumna
    rx = Cm(9.5)
    add_text(sl, "Zatwierdzenie raportu", rx, Cm(3.55), Cm(4.5), Cm(0.5),
             size=13, bold=True, color=NAVY)
    feat_row(sl, "✅","Ofertowy / Bez oferty / Niezatwierdzone","Admin oznacza każdy raport statusem",
             rx, Cm(4.1), Cm(4.5))
    feat_row(sl, "🔓","Odblokuj raport na 24h","Gdy pracownik zgłosi błąd do korekty",
             rx, Cm(5.1), Cm(4.5))

    add_text(sl, "Eksport do XLSX", rx, Cm(6.25), Cm(4.5), Cm(0.5),
             size=13, bold=True, color=NAVY)

    xlsx_box = add_rect(sl, rx, Cm(6.82), Cm(4.5), Cm(1.8),
                        fill=RGBColor(0xff,0xfb,0xeb),
                        line_color=RGBColor(0xfd,0xe6,0x8a), line_w=Pt(1.5))
    xlsx_box.adjustments[0] = 0.1
    add_text(sl, "📥  Arkusz <<Raporty>>", rx+Cm(0.3), Cm(6.88), Cm(4.0), Cm(0.42),
             size=11, bold=True, color=INK)
    xlsx_rows = ["• Data · Pracownik · Ekipa · Godziny",
                 "• Lokalizacja · Wydział · Opis pracy",
                 "• Pojazdy (tablice rejestr.) · Km łącznie",
                 "• Podpisany: Tak / — (sygnatariusze)"]
    for i, row2 in enumerate(xlsx_rows):
        add_text(sl, row2, rx+Cm(0.3), Cm(7.32)+i*Cm(0.3), Cm(4.1), Cm(0.3),
                 size=9, color=INK)

# ──────────────────────────────────────────────────────────────────────────────
slide_cover()
slide_erp()
slide_login()
slide_hr()
slide_rental()
slide_materials()
slide_report()
slide_settlement()

out = "/home/SLEID5/projects/kahma/prezentacja/KAHMA_Szkolenie.pptx"
prs.save(out)
print(f"Saved: {out}")
